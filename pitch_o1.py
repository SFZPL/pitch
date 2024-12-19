import streamlit as st
import google.generativeai as genai
import os
import logging
from datetime import datetime
import re
import time
from typing import List, Optional
from io import BytesIO
import base64
import logging.handlers
import textwrap  # Added import
import json

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import tiktoken

# --------------------------------------------
# Constants & Globals
# --------------------------------------------
ALLOWED_FILE_TYPES = [
    "text/plain",
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
]

MODEL_NAME = "gpt-3.5-turbo-16k"  # keep this for token counting
GEMINI_MODEL_NAME = "gemini-pro"  # Name of the gemini model to use
MAX_CONTEXT = 16384  # max context length for gpt-3.5-turbo-16k
GEMINI_MAX_OUTPUT_TOKENS = 2048  # Max output tokens for gemini-pro

def clean_text(text: str) -> str:
    text = re.sub(r'-\n', '', text)
    text = re.sub(r'\n+', '\n', text)
    text = re.sub(r'[ \t]+', ' ', text)
    text = ''.join(char for char in text if char.isprintable())
    return text.strip()

@st.cache_data
def count_tokens(text: str, model: str = MODEL_NAME) -> int:
    """Counts the number of tokens in a given text using tiktoken."""
    try:
        encoding = tiktoken.encoding_for_model(model)
        return len(encoding.encode(str(text)))
    except Exception as e:
        logging.error(f"Token counting error: {e}")
        return 0

@st.cache_data
def extract_text_from_file(file) -> Optional[str]:
    """Extracts text content from different file types."""
    if file.type not in ALLOWED_FILE_TYPES:
        st.error(f"❌ Unsupported file type: {file.type}")
        return None

    try:
        content = ""
        file_type = file.type

        if file_type == "text/plain":
            content = file.read().decode("utf-8")

        elif file_type == "application/pdf":
            from PyPDF2 import PdfReader
            pdf_reader = PdfReader(file)
            content = "\n".join([page.extract_text() for page in pdf_reader.pages if page.extract_text()])

        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            from docx import Document
            doc = Document(file)
            content = "\n".join([para.text for para in doc.paragraphs if para.text])

        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            from pptx import Presentation
            ppt = Presentation(file)
            content = "\n".join([
                paragraph.text
                for slide in ppt.slides
                for shape in slide.shapes
                if shape.has_text_frame
                for paragraph in shape.text_frame.paragraphs
            ])

        return clean_text(content)

    except Exception as e:
        logging.error(f"File extraction error: {e}")
        st.error(f"❌ Error processing file: {e}")
        return None

def dynamic_margin(prompt_tokens: int) -> int:
    """Calculates a dynamic margin based on the prompt size."""
    # Margin is proportional to the prompt size, min 100, max 500, ~5% of prompt_tokens
    return max(100, min(500, int(prompt_tokens * 0.05)))

def compute_max_tokens(messages: List[dict], desired_completion_tokens: int) -> int:
    """Computes the maximum tokens available for a completion."""
    prompt_tokens = sum(count_tokens(m["content"], model=MODEL_NAME) for m in messages)
    margin = dynamic_margin(prompt_tokens)
    if prompt_tokens + desired_completion_tokens > MAX_CONTEXT:
        available = MAX_CONTEXT - prompt_tokens - margin
        return max(500, min(desired_completion_tokens, available))
    return desired_completion_tokens

def setup_logging():
    """Configures logging to both console and file."""
    try:
        log_formatter = logging.Formatter('%(asctime)s - %(levelname)s: %(message)s')
        log_file = f'pitch_deck_logs_{datetime.now().strftime("%Y%m%d")}.log'

        # Rotating File Handler
        file_handler = logging.handlers.RotatingFileHandler(
            log_file,
            maxBytes=10 * 1024 * 1024,  # 10MB
            backupCount=5
        )
        file_handler.setFormatter(log_formatter)

        # Console Handler
        console_handler = logging.StreamHandler()
        console_handler.setFormatter(log_formatter)

        # Configure root logger
        logger = logging.getLogger()
        logger.addHandler(file_handler)
        logger.addHandler(console_handler)
        logger.setLevel(logging.INFO)

        logging.info("Logging successfully configured.")
    except Exception as e:
        print(f"Error setting up logging: {e}")
        raise

setup_logging()
try:
    # Access the API key from st.secrets
    genai.configure(api_key=st.secrets["google"]["api_key"])
    model = genai.GenerativeModel(GEMINI_MODEL_NAME)
    logging.info("Google API key validated successfully.")
except Exception as e:
    logging.error(f"Google API initialization error: {e}")
    st.error("❌ Google API initialization error.")
    raise

def summarize_content(content: str, prompt_tokens: int = 0) -> Optional[str]:
    """Summarizes content using Gemini."""
    margin = dynamic_margin(prompt_tokens) if prompt_tokens else 500
    desired_completion_tokens = 1500  # shorter summary
    messages = [
        {"role": "system", "content": "You are a helpful assistant that summarizes text."},
        {"role": "user", "content": f"Please provide a concise summary of the following content:\n{content}"}
    ]
    prompt_size = sum(count_tokens(m["content"], model=MODEL_NAME) for m in messages)  # keep this for tokens as gemini doesn't count tokens
    available = MAX_CONTEXT - prompt_size - margin
    max_output_tokens = max(500, min(desired_completion_tokens, available))

    try:
        response = model.generate_content(
            contents=[f"""
        You are a helpful assistant that summarizes text.
        Please provide a concise summary of the following content:\n{content}
        """],
            generation_config=genai.types.GenerationConfig(max_output_tokens=max_output_tokens, temperature=0.5)
        )
        summary = response.text
        return summary.strip() if summary else None
    except Exception as e:
        logging.error(f"Content summarization error: {e}")
        st.error(f"❌ Error summarizing content: {e}")
        return None

class DocumentGenerator:
    """A class to generate and manipulate documents (pitch decks, corporate profiles)."""
    def __init__(self):
        self.pitch_deck_prompt = """
You are an expert pitch deck content generator trained in creating world-class business presentations.
Your goal is to transform raw content into a compelling narrative that captures investor attention.

Key Presentation Principles:
1. Tell a Story: Create a clear, engaging narrative
2. Highlight Problem and Solution
3. Demonstrate Market Opportunity
4. Showcase Unique Value Proposition
5. Use Data-Driven Insights
6. Maintain Professional and Exciting Tone

Follow this Strict Structure:
- Title Slide: Company Name, Tagline
- Problem Statement: Clear, Impactful
- Solution: Innovative Approach
- Market Analysis: Size, Growth, Opportunity
- Business Model
- Competitive Landscape
- Financial Projections
- Team Overview
- Call to Action

Output Format:
For each slide, provide a JSON object in the following structure:
{
    "slide_title": "[Slide Title Here]",
    "content": "[Concise slide content here]",
    "suggested_visual": "[Description or URL for a suitable visual]"
}
Ensure your JSON output is always valid.
"""

        self.corporate_profile_prompt = """
You are a professional corporate profile content generator. Your task is to create a comprehensive and engaging corporate profile presentation (PowerPoint style) that effectively represents the company's brand and offering.

Key Elements to Include:
1. Title Slide: Company Name, Logo, and Tagline
2. Mission, Vision, and Values
3. Company History: Founding story, key milestones
4. Products/Services: Detailed descriptions, unique selling points
5. Market Position: Industry standing, competitive advantages
6. Achievements & Awards: Recognitions, milestones
7. Team & Leadership Overview
8. Future Goals & Strategic Plans
9. Contact Information

Output Format:
For each slide, provide a JSON object in the following structure:
{
    "slide_title": "[Slide Title Here]",
    "content": "[Concise slide content here]",
    "suggested_visual": "[Description or URL for a suitable visual]"
}
Ensure your JSON output is always valid.
"""

        self.token_cache = {}  # Initialize token cache

    def sanitize_input(self, text: str) -> str:
        """Removes potentially malicious characters from input text."""
        return re.sub(r'[<>&\'"]', '', text)

    def hex_to_rgb(self, hex_color):
        """Converts a hex color code to an RGB tuple."""
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def export_to_pptx(self, slides: List[dict]) -> Optional[BytesIO]:
        """Exports the slides data to a PowerPoint presentation."""
        try:
            prs = Presentation()
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)

            colors = {
                'background': 'F0F4F8',
                'primary': '1E40AF',
                'secondary': '3B82F6',
                'accent': 'EF4444',
            }

            for slide_data in slides:
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(colors['background']))

                slide_title = slide_data.get('slide_title', "")
                slide_body = slide_data.get('content', "")
                # slide_visual = slide_data.get('suggested_visual', "")  # Not used

                if slide_title:
                    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1.5))
                    title_frame = title_box.text_frame
                    title_frame.text = slide_title
                    title_frame.paragraphs[0].font.size = Pt(44)
                    title_frame.paragraphs[0].font.color.rgb = RGBColor(*self.hex_to_rgb(colors['primary']))
                    title_frame.paragraphs[0].font.bold = True

                if slide_body:
                    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(15), Inches(6))
                    content_frame = content_box.text_frame
                    content_frame.text = slide_body
                    content_frame.paragraphs[0].font.size = Pt(24)
                    content_frame.paragraphs[0].font.color.rgb = RGBColor(*self.hex_to_rgb(colors['secondary']))

                # Optionally, incorporate the visual suggestion into the PPTX.
                # Currently, it's only represented textually in the interface.
                # No changes here unless you want to add a placeholder shape or something similar.

            pptx_file = BytesIO()
            prs.save(pptx_file)
            pptx_file.seek(0)
            return pptx_file

        except Exception as e:
            logging.error(f"PowerPoint export error: {e}")
            st.error(f"❌ Failed to export PowerPoint: {e}")
            return None

    def _count_tokens(self, text: str, model: str = MODEL_NAME) -> int:
        """Counts tokens for a string using the cache."""
        cache_key = (text, model)
        if cache_key not in self.token_cache:
            self.token_cache[cache_key] = count_tokens(text, model)
        return self.token_cache[cache_key]

    def generate_document(
            self,
            content: str,
            document_type: str = "Pitch Deck",
            max_retries: int = 3
        ) -> Optional[List[dict]]:
        """Generates a document (pitch deck or corporate profile) using Gemini."""

        # Choose the built-in base system prompt
        if document_type == "Pitch Deck":
            system_prompt = self.pitch_deck_prompt
        elif document_type == "Corporate Profile":
            system_prompt = self.corporate_profile_prompt
        else:
            st.error("❌ Unsupported document type selected.")
            return None

        for attempt in range(max_retries):
            try:
                # Create the base prompt
                full_prompt = f"""
                {system_prompt}

                INSTRUCTIONS:
                - Split the content into multiple slides.
                - Each slide must include:
                1. A title under 'slide_title'.
                2. Concise content under 'content'.
                3. A suggested visual under 'suggested_visual'.
                - Ensure the output follows this JSON format:
                {{
                    "slide_title": "Title of Slide",
                    "content": "Concise slide content here.",
                    "suggested_visual": "Description or URL for a suitable visual."
                }}
                """

                # Append the content to be processed
                full_prompt += f"\nCONTENT TO PROCESS:\n{content}\n"

                # Call the AI model
                response = model.generate_content(
                    contents=[full_prompt],
                    generation_config=genai.types.GenerationConfig(
                        max_output_tokens=GEMINI_MAX_OUTPUT_TOKENS, 
                        temperature=0.7
                    )
                )

                # Log the raw output for debugging
                generated_content = response.text
                logging.debug(f"Raw Generated Content: {generated_content}")

                # Extract multiple JSON objects
                slides = []
                potential_json_strings = re.findall(r'\{.*?\}', generated_content, re.DOTALL)
                for json_str in potential_json_strings:
                    try:
                        slide = json.loads(json_str)
                        if "slide_title" in slide and "content" in slide:
                            slides.append(slide)
                    except json.JSONDecodeError:
                        logging.error(f"JSON decode error for: {json_str}")
                        continue

                # Fallback: If only one slide is generated, attempt to reprocess
                if len(slides) < 2:
                    logging.warning("Generated content included insufficient slides. Retrying...")
                    continue  # Retry if slides are insufficient

                # Validate if multiple slides are generated
                if len(slides) > 1:
                    logging.info(f"Document generated successfully with {len(slides)} slides.")
                    return slides

            except Exception as e:
                logging.error(f"Error during document generation: {e}")
                if attempt < max_retries - 1:
                    time.sleep(2)  # Backoff before retrying
                else:
                    st.error("❌ Document generation failed after multiple attempts.")
                    return None


    def update_section(
        self,
        section_number: int,
        section_content: dict,
        edit_instructions: str,
        previous_sections: List[dict],
        document_type: str
    ) -> Optional[dict]:
        """Updates a specific section of the document based on user instructions."""
        edit_instructions = self.sanitize_input(edit_instructions)

        if document_type == "Pitch Deck":
            system_prompt = self.pitch_deck_prompt
        elif document_type == "Corporate Profile":
            system_prompt = self.corporate_profile_prompt
        else:
            st.error("❌ Unsupported document type selected.")
            return None

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": f"You are to update Slide {section_number} of the {document_type}."},
            {"role": "user", "content": f"Previous Slides:\n" + json.dumps(previous_sections)},
            {"role": "user", "content": f"Current Slide {section_number} Content:\n{json.dumps(section_content)}"},
            {"role": "user", "content": f"Edit Instructions:\n{edit_instructions}"}
        ]

        prompt_tokens = sum(self._count_tokens(m["content"], model=MODEL_NAME) for m in messages)
        margin = dynamic_margin(prompt_tokens)
        desired_completion_tokens = max(1000, MAX_CONTEXT - prompt_tokens - margin)
        max_tokens = max(500, min(desired_completion_tokens, MAX_CONTEXT - prompt_tokens - margin))

        try:
            response = model.generate_content(
                contents=[f"""
                {system_prompt}
                You are to update Slide {section_number} of the {document_type}.
                Previous Slides: {json.dumps(previous_sections)}
                Current Slide {section_number} Content: {json.dumps(section_content)}
                Edit Instructions: {edit_instructions}
                """],
                generation_config=genai.types.GenerationConfig(max_output_tokens=GEMINI_MAX_OUTPUT_TOKENS, temperature=0.7)
            )
            updated_section_content = response.text
            try:
                updated_section = json.loads(updated_section_content)
                logging.info(f"Slide {section_number} updated successfully.")
                return updated_section
            except json.JSONDecodeError as e:
                logging.error(f"JSON decode error during slide update: {e}")
                st.error(f"❌ Error decoding JSON response for slide update. Details: {e}")
                return None

        except Exception as e:
            logging.error(f"Slide update error: {e}")
            st.error(f"❌ Error updating slide: {e}")
            return None

    def analyze_existing_presentation(self, content: str, document_type: str) -> Optional[str]:
        """Analyzes an existing presentation and provides feedback."""
        analysis_prompt = f"""
You are a presentation analyst. The user has provided an existing {document_type}. 
Please analyze the content and provide:
- Key strengths: Which aspects are well done?
- Potential weaknesses or areas for improvement: What could be improved?
- Suggested tweaks or enhancements: How to strengthen the narrative, visuals, or clarity?

Keep the tone constructive and professional.
"""
        content = self.sanitize_input(content)
        messages = [
            {"role": "system", "content": analysis_prompt},
            {"role": "user", "content": f"Existing {document_type} Content:\n{content}"}
        ]

        prompt_tokens = sum(count_tokens(m["content"], model=MODEL_NAME) for m in messages)
        margin = dynamic_margin(prompt_tokens)
        desired_completion_tokens = max(1000, MAX_CONTEXT - prompt_tokens - margin)
        max_tokens = max(500, min(desired_completion_tokens, MAX_CONTEXT - prompt_tokens - margin))

        try:
            response = model.generate_content(
                contents=[f"""
                {analysis_prompt}
                Existing {document_type} Content: {content}
                """],
                generation_config=genai.types.GenerationConfig(max_output_tokens=GEMINI_MAX_OUTPUT_TOKENS, temperature=0.7)
            )
            analysis = response.text
            return analysis.strip()
        except Exception as e:
            logging.error(f"Presentation analysis error: {e}")
            st.error(f"❌ Error analyzing presentation: {e}")
            return None

    def apply_global_edit(self, instruction: str, slides: List[dict], document_type: str, max_retries: int = 3) -> Optional[List[dict]]:
        """Applies a global edit instruction to all slides."""
        edited_slides = []
        for i, slide in enumerate(slides, 1):
            for attempt in range(max_retries):
                try:
                    system_prompt = self.pitch_deck_prompt if document_type == "Pitch Deck" else self.corporate_profile_prompt

                    full_prompt = f"""
                    {system_prompt}

                    INSTRUCTIONS:
                    - Apply the following edit to the slide.
                    - Ensure the slide remains coherent and follows the original structure.

                    EDIT INSTRUCTION:
                    {instruction}

                    CURRENT SLIDE CONTENT:
                    {{
                        "slide_title": "{slide.get('slide_title', '')}",
                        "content": "{slide.get('content', '')}",
                        "suggested_visual": "{slide.get('suggested_visual', '')}"
                    }}

                    UPDATED SLIDE CONTENT:
                    {{
                        "slide_title": "[Slide Title Here]",
                        "content": "[Updated slide content here]",
                        "suggested_visual": "[Updated description or URL for a suitable visual]"
                    }}
                    Ensure your JSON output is always valid.
                    """

                    response = model.generate_content(
                        contents=[full_prompt],
                        generation_config=genai.types.GenerationConfig(
                            max_output_tokens=GEMINI_MAX_OUTPUT_TOKENS, 
                            temperature=0.7
                        )
                    )

                    updated_content = response.text
                    updated_slide = json.loads(updated_content)

                    if all(key in updated_slide for key in ("slide_title", "content", "suggested_visual")):
                        edited_slides.append(updated_slide)
                        logging.info(f"Slide {i} edited successfully.")
                        break  # Exit retry loop on success
                    else:
                        logging.error(f"Missing keys in updated slide {i}. Retrying...")
                        continue

                except json.JSONDecodeError as e:
                    logging.error(f"JSON decode error for slide {i}: {e}")
                except Exception as e:
                    logging.error(f"Error applying global edit to slide {i}: {e}")

                if attempt < max_retries - 1:
                    time.sleep(2)  # Backoff before retrying
                else:
                    st.error(f"❌ Failed to apply global edit to slide {i} after multiple attempts.")
                    return None

        return edited_slides

def format_full_document_view(sections):
    """Formats the generated document sections into a full, readable document view."""
    full_document = []
    for i, section_data in enumerate(sections, 1):
        slide_title = section_data.get('slide_title', f"Slide {i}")
        slide_body = section_data.get('content', "")
        slide_visual = section_data.get('suggested_visual', "")

        section_text = f"**Slide {i}:**\n\n" \
                       f"**Title:** {slide_title}\n\n" \
                       f"**Content:**\n\n{slide_body}\n\n" \
                       f"**Suggested visual idea:**\n\n{slide_visual}\n\n"
        full_document.append(section_text)

    return "\n\n".join(full_document)

# --------------------------------------------
# Main Application
# --------------------------------------------

def main():
    st.set_page_config(page_title="Prez.AI Document Generator", page_icon="📊", layout="wide")

    st.markdown("""
    <style>
    .main-header {
        font-size: 2.5em;
        color: #333333;
        text-align: center;
        margin-bottom: 20px;
        font-weight: bold;
    }
    .subheader {
        color: #666666;
        text-align: center;
        margin-bottom: 30px;
        font-size: 1.2em;
    }
    .custom-text {
        background-color: #f9f9f9;
        border-radius: 5px;
        padding: 15px;
        border: 1px solid #e0e0e0;
    }
    .stButton>button {
        background-color: #4CAF50;
        color: white;
        font-weight: bold;
        border: none;
        padding: 10px 20px;
        font-size: 16px;
        margin: 4px 2px;
        transition-duration: 0.4s;
        cursor: pointer;
    }
    .stButton>button:hover {
        background-color: #45a049;
    }
    .chat-message {
        margin: 10px 0;
        padding: 10px;
        border-radius: 5px;
        font-size: 14px;
    }
        .user { background-color: #f0f0f0; }
    .assistant { background-color: #e6f3ff; }
    .system { background-color: #e6ffe6; }
    .error { background-color: #ffe6e6; }
    .copy-area {
    border: 1px solid #ccc;
    padding: 10px;
    border-radius: 5px;
    margin-top: 10px;
    background-color: #f9f9f9;
    }
    .tooltip {
        position: relative;
        display: inline-block;
    }

    .tooltip .tooltiptext {
        visibility: hidden;
        width: 200px;
        background-color: #555;
        color: #fff;
        text-align: center;
        border-radius: 6px;
        padding: 5px;
        position: absolute;
        z-index: 1;
        bottom: 125%;
        left: 50%;
        margin-left: -100px;
        opacity: 0;
        transition: opacity 0.3s;
    }
    .tooltip:hover .tooltiptext {
        visibility: visible;
        opacity: 1;
    }
    </style>
    """, unsafe_allow_html=True)

    st.markdown('<div class="main-header">Prez.AI Document Generator 🎨📈</div>', unsafe_allow_html=True)
    st.markdown('<div class="subheader">Transform Your Ideas into Compelling Documents 🚀</div>', unsafe_allow_html=True)

    if 'doc_generator' not in st.session_state:
        st.session_state.doc_generator = DocumentGenerator()

    session_defaults = {
        'uploaded_file': None,
        'generated_document_sections': [],
        'chat_history': [],
        'current_step': 'initial',
        'file_content': None,
        'document_type': None,
        'analysis_result': None,
        'view_mode': "Slide View",
        'action': None
    }

    for key, default_value in session_defaults.items():
        if key not in st.session_state:
            st.session_state[key] = default_value

    with st.container():

        uploaded_file = st.file_uploader(
            "📂 Upload Your Source Document",
            type=["txt", "pdf", "docx", "pptx"],
            help="Supported formats: Text, PDF, Word, PowerPoint"
        )


        if uploaded_file != st.session_state.uploaded_file:
            st.session_state.uploaded_file = uploaded_file
            st.session_state.generated_document_sections = []
            st.session_state.chat_history = []
            st.session_state.file_content = None
            st.session_state.analysis_result = None
            st.session_state.current_step = 'file_uploaded'
            st.session_state.document_type = None
            st.session_state.action = None

        if uploaded_file:
            st.markdown("### 📝 Processing Your File")

            # Display a loading spinner and extract file content
            if not st.session_state.file_content:
                with st.spinner("🕒 Extracting file content... Please wait!"):
                    st.session_state.file_content = extract_text_from_file(uploaded_file)

                if st.session_state.file_content:
                    st.markdown("""
                    <div style="background-color: #e6f4ea; padding: 10px; border-radius: 5px; margin: 10px 0;">
                        <p style="color: #2e7d32; font-weight: bold; font-size: 16px;">
                        ✅ File content extracted successfully! Your document is ready for processing.
                        </p>
                    </div>
                    """, unsafe_allow_html=True)
                else:
                    st.markdown("""
                    <div style="background-color: #ffebee; padding: 10px; border-radius: 5px; margin: 10px 0;">
                        <p style="color: #b71c1c; font-weight: bold; font-size: 16px;">
                        ❌ Could not extract content from the file. Please upload a supported file type.
                        </p>
                    </div>
                    """, unsafe_allow_html=True)

            # Display next steps only after successful extraction
            if st.session_state.file_content:
                st.markdown("### 🛠️ Next Step")
                
                # Select action
                action = st.selectbox(
                    "Select action:",
                    ["Generate Pitch Deck", "Generate Corporate Profile", "Analyze Presentation"],
                    index=0
                )
                st.session_state.action = action

                # Confirm button
                if st.button("Confirm"):
                    if st.session_state.action == "Generate Pitch Deck":
                        with st.spinner("🤖 Generating your pitch deck... Please wait!"):
                            doc = st.session_state.doc_generator.generate_document(
                                content=st.session_state.file_content,
                                document_type="Pitch Deck"
                            )
                            if doc:
                                st.success("✅ **Pitch deck generated successfully!**")
                                st.session_state.generated_document_sections = doc
                                st.session_state.document_type = "Pitch Deck"
                            else:
                                st.error("❌ Failed to generate pitch deck.")
                    elif st.session_state.action == "Generate Corporate Profile":
                        with st.spinner("🤖 Generating your corporate profile... Please wait!"):
                            doc = st.session_state.doc_generator.generate_document(
                                content=st.session_state.file_content,
                                document_type="Corporate Profile"
                            )
                            if doc:
                                st.success("✅ **Corporate profile generated successfully!**")
                                st.session_state.generated_document_sections = doc
                                st.session_state.document_type = "Corporate Profile"
                            else:
                                st.error("❌ Failed to generate corporate profile.")
                    elif st.session_state.action == "Analyze Presentation":
                        if uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                            with st.spinner("🔍 Analyzing presentation..."):
                                analysis = st.session_state.doc_generator.analyze_existing_presentation(
                                    content=st.session_state.file_content,
                                    document_type="Presentation"
                                )
                                if analysis:
                                    st.success("✅ **Presentation analysis complete!**")
                                    st.session_state.analysis_result = analysis
                                else:
                                    st.error("❌ Failed to analyze the presentation.")
                        else:
                            st.error("❌ Presentation analysis is only supported for PPTX files.")

                    else:
                         st.error("❌ Please select an action.")

        else:
            st.info("📄 Upload a document to begin.")

        if st.session_state.generated_document_sections:

            view_mode = st.radio(
                "View Mode",
                ["Slide View", "Full Document View"],
                 horizontal=True,
                 help="Choose how you want to view the generated document",
                 index=["Slide View", "Full Document View"].index(st.session_state.view_mode)
            )

            st.session_state.view_mode = view_mode

            if st.session_state.view_mode == "Slide View":
                slide_tabs = st.tabs([f"Slide {i+1}" for i in range(len(st.session_state.generated_document_sections))])

                for i, section_data in enumerate(st.session_state.generated_document_sections):
                     with slide_tabs[i]:
                        slide_title = section_data.get('slide_title', 'No Title')
                        slide_body = section_data.get('content', '')
                        slide_visual = section_data.get('suggested_visual', '')
                        
                        st.subheader(slide_title)

                        st.markdown("**Content:**")
                        if slide_body:
                            st.markdown(slide_body, unsafe_allow_html=True)
                        else:
                            st.markdown("_No content provided._")
                        
                        st.markdown("**Suggested Visual Idea:**")
                        if slide_visual:
                            if re.match(r'^https?://', slide_visual):
                                st.image(slide_visual)
                            else:
                                # If the visual is not a URL, display as placeholder
                                st.markdown(f"{slide_visual}")
                        else:
                            st.markdown("_No visual provided._")
                        
                        with st.expander(f"✏️ Edit Slide {i+1}"):
                            
                            user_input = st.text_input(f"Edit Slide {i + 1} Content:", key=f"slide_edit_{i}", placeholder="Enter your instructions for this slide...")
                            if st.button(f"Submit Edit Slide {i+1}", key=f"submit_edit_{i}"):
                                if user_input:
                                    st.session_state.chat_history.append({
                                        'type': 'user',
                                        'message': f"✍️ Edit Slide {i+1}: {user_input}"
                                    })
                                    with st.spinner("🔧 Updating slide..."):
                                        updated_section = st.session_state.doc_generator.update_section(
                                            section_number=i+1,
                                            section_content=section_data,
                                            edit_instructions=user_input,
                                            previous_sections=st.session_state.generated_document_sections,
                                            document_type=st.session_state.document_type
                                        )
                                        if updated_section:
                                            st.session_state.generated_document_sections[i] = updated_section
                                            st.session_state.chat_history.append({
                                                'type': 'assistant',
                                                'message': f"✅ Updated Slide {i+1}:\n{json.dumps(updated_section)}"
                                            })
                                            st.success(f"Slide {i+1} updated successfully!")
                                            st.rerun()
                                        else:
                                            st.error("❌ Failed to update the slide.")
                                else:
                                    st.warning("⚠️ Please provide edit instructions.")

                # Add the option to edit all slides at once
                st.markdown("---")
                st.header("📝 Apply Global Edit to All Slides")

                with st.expander("✏️ Global Edit"):
                    global_edit_instruction = st.text_input(
                        "Enter your global edit instruction:",
                        placeholder="e.g., 'Add more details to all slides', 'Make the tone more formal'"
                    )
                    if st.button("Apply Global Edit"):
                        if global_edit_instruction.strip():
                            st.session_state.chat_history.append({
                                'type': 'user',
                                'message': f"✍️ Global Edit: {global_edit_instruction.strip()}"
                            })
                            with st.spinner("🔧 Applying global edit to all slides..."):
                                updated_slides = st.session_state.doc_generator.apply_global_edit(
                                    instruction=global_edit_instruction.strip(),
                                    slides=st.session_state.generated_document_sections,
                                    document_type=st.session_state.document_type
                                )
                                if updated_slides:
                                    st.session_state.generated_document_sections = updated_slides
                                    st.session_state.chat_history.append({
                                        'type': 'assistant',
                                        'message': "✅ Global edit applied successfully to all slides!"
                                    })
                                    st.success("✅ Global edit applied successfully to all slides!")
                                    st.rerun()
                                else:
                                    st.error("❌ Failed to apply global edit to all slides.")
                        else:
                            st.warning("⚠️ Please enter a valid edit instruction.")

            else:
                 st.header("📄 Full Document View")
                
                 # Format the full document
                 full_document_text = format_full_document_view(st.session_state.generated_document_sections)
                
                # Display the full document
                 st.markdown(full_document_text)

            col2_1, col2_2 = st.columns(2)
            with col2_1:
                if st.session_state.view_mode == "Slide View":
                   # Download slides as JSON
                    slide_texts = [json.dumps(slide) for slide in st.session_state.generated_document_sections]
                    st.download_button(
                        label="💾 Download Slides as JSON",
                        data="\n\n".join(slide_texts),
                        file_name=f"{st.session_state.document_type.lower().replace(' ', '_')}.json",
                        mime="application/json"
                    )
                else:
                     # Download full document view as text
                    st.download_button(
                        label="💾 Download Full Document as Text",
                        data=format_full_document_view(st.session_state.generated_document_sections),
                        file_name=f"{st.session_state.document_type.lower().replace(' ', '_')}.txt",
                        mime="text/plain"
                    )

            with col2_2:
                export_pptx = st.button("📂 Export as PowerPoint")
                if export_pptx:
                    with st.spinner("Exporting to PowerPoint..."):
                        pptx_file = st.session_state.doc_generator.export_to_pptx(
                            slides=st.session_state.generated_document_sections
                        )
                        if pptx_file:
                            b64 = base64.b64encode(pptx_file.getvalue()).decode()
                            href = f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}" download="{st.session_state.document_type.lower().replace(" ", "_")}.pptx">📥 Download PowerPoint</a>'
                            st.markdown(href, unsafe_allow_html=True)
                        else:
                            st.error("❌ Failed to export PowerPoint.")

        if st.session_state.analysis_result:
            st.header("🔍 Presentation Analysis")
            st.markdown(st.session_state.analysis_result, unsafe_allow_html=True)

    st.markdown("---")
    st.subheader("📜 Session History")
    chat_container = st.container()

    with chat_container:
        for chat_entry in st.session_state.chat_history:
            message_class = chat_entry.get('type', 'user')
            st.markdown(f'<div class="chat-message {message_class}">{chat_entry["message"]}</div>', unsafe_allow_html=True)

if __name__ == "__main__":
    main()
